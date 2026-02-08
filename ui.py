import streamlit as st
import requests
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------------- PAGE CONFIG ----------------
st.set_page_config(
    page_title="HPCL Lead Intelligence",
    layout="centered"
)

# ---------------- HEADER ----------------
st.title("HPCL B2B Lead Intelligence")
st.caption("Mobile Sales Officer View")
st.divider()

# ---------------- SESSION STATE ----------------
if "lead_generated" not in st.session_state:
    st.session_state.lead_generated = False

if "notify_clicked" not in st.session_state:
    st.session_state.notify_clicked = False

# ---------------- INPUT SECTION ----------------
st.subheader("New Lead Discovery")

with st.container():
    company = st.text_input("Company Name")
    location = st.text_input("Location")

    text = st.text_area(
        "Public Signal (Tender / News / Website)",
        height=140,
        placeholder="Tender for supply of Furnace Oil for captive power plant..."
    )

# ---------------- GENERATE LEAD ----------------
if st.button("Generate Lead Dossier", use_container_width=True):
    if not company or not location or not text:
        st.warning("Please fill all fields")
    else:
        response = requests.post(
            "http://127.0.0.1:8000/analyze",
            json={
                "company": company,
                "location": location,
                "text": text
            }
        )
        st.session_state.lead_data = response.json()
        st.session_state.lead_generated = True
        st.session_state.notify_clicked = False

# ---------------- SHOW LEAD DOSSIER ----------------
if st.session_state.lead_generated:
    data = st.session_state.lead_data

    st.divider()
    st.subheader("Lead Dossier")

    # ---- BASIC INFO CARD ----
    with st.container():
        st.markdown(f"""
        **Company:** {data['company']}  
        **Location:** {data['location']}  
        **Intent Type:** {data['intent_type']}
        """)

    # ---- METRICS ----
    col1, col2 = st.columns(2)
    with col1:
        st.metric("Urgency", data["urgency"])
    with col2:
        st.metric("Confidence Score", data["confidence_score"])

    # ---- PRODUCTS ----
    st.markdown("**Recommended Products**")
    for p in data["recommended_products"]:
        st.write(f"- {p}")

    # ---- REASONS ----
    st.markdown("**Reason Codes**")
    for r in data["reason_codes"]:
        st.write(f"- {r}")

    st.success(data["suggested_next_action"])

    # ---------------- ACTIONS ----------------
    st.divider()
    st.subheader("Sales Actions")

    col1, col2 = st.columns(2)

    # ---- NOTIFY ----
    with col1:
        if st.button("Notify Sales Officer", use_container_width=True):
            st.session_state.notify_clicked = True

        if st.session_state.notify_clicked:
            st.success("Sales alert triggered (sandbox)")

    # ---- PPT EXPORT ----
    with col2:
        if st.button("Download Lead PPT", use_container_width=True):
            prs = Presentation()

            # Slide 1 – Overview
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "HPCL Lead Dossier"
            content = slide.placeholders[1]
            content.text = (
                f"Company: {data['company']}\n"
                f"Location: {data['location']}\n"
                f"Intent: {data['intent_type']}\n"
                f"Urgency: {data['urgency']}\n"
                f"Confidence: {data['confidence_score']}"
            )

            # Slide 2 – Products & Reasons
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Products & Signals"
            content = slide.placeholders[1]
            content.text = (
                "Recommended Products:\n"
                + ", ".join(data["recommended_products"])
                + "\n\nReason Codes:\n"
                + ", ".join(data["reason_codes"])
            )

            # Slide 3 – Action
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Next Action"
            slide.placeholders[1].text = data["suggested_next_action"]

            ppt_path = "lead_dossier.pptx"
            prs.save(ppt_path)

            with open(ppt_path, "rb") as f:
                st.download_button(
                    label="Download PPT",
                    data=f,
                    file_name="HPCL_Lead_Dossier.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
